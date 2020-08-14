from pptx import Presentation
from pptx.util import Pt,Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# HELPER FUNCTIONS ----------------------------------------------------------------------------------------------------
def numberdomains(project):
    return len(project["workitems"].keys())


def findslide(target):
    for s in deck.slides:
        if s.shapes.title != None and target in s.shapes.title.text:
            return s

def insert(tab,r,c,txt):
    tab.cell(r,c).text = txt
    if r == 0:
        tab.cell(r,c).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

def setup(tab):
    for r in range(len(tab.rows)):
        for c in range(len(tab.columns)):
            tab.cell(r,c).fill.solid()
            tab.cell(r,c).fill.fore_color.rgb = RGBColor(255,255,255)
            tab.cell(r,c).text_frame.paragraphs[0].font.size = Pt(15)
            tab.cell(r,c).text_frame.paragraphs[0].font.color.rgb = RGBColor(49,56,64)

outputf = '/tmp/new.pptx'
deck = None
# Add template name here
def generateproposal(project, inputf):
    global deck
    deck = Presentation(inputf)
    # ------------------------------------ TITLE SLIDE ------------------------------------

    # Getting and adding title slide object
    slide = findslide("Internal Control Review Proposal")

    # Inserting content into title of title slide
    text = slide.shapes.title.text + "\n\n %s" %(project['name'])
    slide.shapes.title.text = text

    # Inserting date to title slide
    slide.placeholders[10].text = '%s' %(project['date'])


    # ------------------------------------ OBJECTIVES SLIDE ------------------------------------
    slide = findslide('Your Objectives')

    # Sub-title
    text = "We understand that %s (“the Company”) would like to engage an experienced service provider with proven track records in the industry as your consultant to achieve the following objectives:" %(project['name'])
    slide.shapes[1].text = text

    # Objectives bullet points (TODO: Formatting for bullets, bullet icon)
    text_frame = slide.shapes[2].text_frame
    text_frame.clear()


    # Inserting objectives into textbox as bullet points
    for obj in project['objectives']:
        para = text_frame.add_paragraph()
        para.text = obj

    # ------------------------------------ PROPOSED FEES SLIDE ------------------------------------
    slide = findslide('Proposed Fee')

    # Adding costs to slide
    slide.placeholders[10].text = str(project['cost1'])
    slide.placeholders[11].text = str(project['cost2'])

    # Adding billing information
    text_frame = slide.placeholders[1].text_frame
    text_frame.clear()
    p = text_frame.paragraphs[0]
    p.text = "Billings:"

    for para_str in project['billings']:
        p = text_frame.add_paragraph()
        p.text = para_str
        p.level = 1

    # ------------------------------------ PROPOSED REVIEW AREAS SLIDES ------------------------------------

    for dom in project["workitems"].keys():
        for slide in deck.slides:
            if slide.shapes.title != None and 'Proposed Review Areas and Coverage' in slide.shapes.title.text:
                for shape in slide.shapes:
                    if shape.has_text_frame and len(shape.text_frame.paragraphs) == 1 and shape.text == dom:
                        for s in slide.shapes:
                            if s.has_text_frame and len(s.text_frame.paragraphs)>5:
                                s.text_frame.clear()
                                for pt in project["workitems"][dom]:
                                    p = s.text_frame.paragraphs[-1]
                                    p.text = pt
                                    s.text_frame.add_paragraph()


    # ------------------------------------ PROPOSED WORK PLAN ------------------------------------
    slide = findslide('Proposed Work Plan')
    if slide != None:
        sh = slide.shapes.add_table(numberdomains(project)+2,3,Inches(0.6),Inches(2.2),Inches(9),Inches(4))
        tab = sh.table
        setup(tab)
        tab.columns[0].width = Inches(2)
        tab.columns[1].width = Inches(5.2)
        tab.columns[2].width = Inches(2)
        insert(tab,0,0,"Proposed Timing of Review")
        insert(tab,0,1,"Proposed Review Areas")
        insert(tab,0,2,"Proposed Entity")
        insert(tab,1,0,project["startdate"])
        for i,s in enumerate(project["workitems"].keys()):
            insert(tab,i+1,1,s)
        insert(tab,1,2,project["name"])
        insert(tab,numberdomains(project)+1,0,project["enddate"])
        insert(tab,numberdomains(project)+1,1,"Post Implementation Follow-Up Review")
        tab.cell(1,0).merge(tab.cell(numberdomains(project),0))
        tab.cell(1,2).merge(tab.cell(numberdomains(project),2))

    # Saving the finished slides
    deck.save(outputf)
    return outputf

